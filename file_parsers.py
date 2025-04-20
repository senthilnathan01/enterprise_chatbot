# file_parsers.py
"""Functions to parse different file types.
Each parser should return a list of tuples:
[(text_segment_1, metadata_1), (text_segment_2, metadata_2), ...]
where metadata includes at least 'source' and potentially 'page_number' etc.
Also returns the list of found URLs.
"""

import streamlit as st
import io
import fitz # PyMuPDF
import docx
from pptx import Presentation
import pandas as pd
import json
from utils import log_message, find_urls
from image_processor import get_gemini_vision_description_ocr

# --- PDF Parser (Refined for Page Numbers) ---
def parse_pdf(file_content, filename):
    parsed_segments = [] # List of (text_segment, metadata)
    urls = set()
    try:
        doc = fitz.open(stream=file_content, filetype="pdf")
        log_message(f"Processing PDF '{filename}': {len(doc)} pages.", "info")
        for page_num, page in enumerate(doc):
            page_num_actual = page_num + 1 # 1-based page number
            page_text = ""
            page_meta = {"source": filename, "page_number": page_num_actual}
            ocr_text_content = "" # Store OCR result separately

            try:
                page_text = page.get_text("text", sort=True) # Get text sorted by position
                if page_text: # Don't add empty pages unless OCR finds something
                    urls.update(find_urls(page_text))
                    parsed_segments.append((page_text, page_meta.copy()))
                else:
                    log_message(f"No extractable text found on page {page_num_actual} of {filename}. Checking for images/OCR.", "debug")

                # OCR Fallback check (can run even if some text was found, for mixed pages)
                ocr_needed = len(page_text.strip()) < 100 # Heuristic: try OCR if very little text
                if not page_text or ocr_needed:
                    log_message(f"Attempting OCR for page {page_num_actual} in '{filename}'.", "info")
                    try:
                        pix = page.get_pixmap(dpi=150)
                        img_bytes = pix.tobytes("png")
                        ocr_text_content = get_gemini_vision_description_ocr(img_bytes, f"{filename}_page{page_num_actual}_ocr")
                        if ocr_text_content and "No description or text extracted" not in ocr_text_content:
                             ocr_meta = page_meta.copy()
                             ocr_meta["content_type"] = "ocr_page"
                             parsed_segments.append((f"[OCR Result for Page {page_num_actual}]\n{ocr_text_content}", ocr_meta))
                        else:
                             log_message(f"OCR found no significant text on page {page_num_actual}.", "debug")
                    except Exception as ocr_err:
                        log_message(f"Error during OCR fallback for {filename} page {page_num_actual}: {ocr_err}", "warning")

                # Extract Embedded Images (if any)
                try:
                    image_list = page.get_images(full=True)
                    for img_index, img_info in enumerate(image_list):
                        xref = img_info[0]
                        base_image = doc.extract_image(xref)
                        if base_image:
                            image_bytes = base_image["image"]
                            img_meta_filename = f"{filename}_page{page_num_actual}_img{img_index}"
                            img_desc_ocr = get_gemini_vision_description_ocr(image_bytes, img_meta_filename)
                            if img_desc_ocr and "No description or text extracted" not in img_desc_ocr:
                                 img_meta = page_meta.copy()
                                 img_meta["content_type"] = "embedded_image"
                                 img_meta["image_ref"] = img_index
                                 parsed_segments.append((f"[Analysis of Embedded Image {img_index} on Page {page_num_actual}]\n{img_desc_ocr}", img_meta))
                        else:
                            log_message(f"Could not extract image xref {xref} on page {page_num_actual}.", "warning")
                except Exception as img_err:
                    log_message(f"Error processing embedded images on page {page_num_actual}: {img_err}", "warning")

            except Exception as page_err:
                 log_message(f"Error processing page {page_num_actual} of {filename}: {page_err}", "warning")

        log_message(f"Finished processing PDF '{filename}'. Found URLs: {len(urls)}")
    except Exception as e:
        log_message(f"Critical error processing PDF '{filename}': {e}", "error")
        # Return empty lists on critical failure
        return [], []
    return parsed_segments, list(urls)

# --- DOCX Parser ---
# Returns list [(full_text, metadata), (image_analysis_1, metadata), ...]
def parse_docx(file_content, filename):
    parsed_segments = []
    urls = set()
    base_meta = {"source": filename, "content_type": "text"}
    img_meta = {"source": filename, "content_type": "embedded_image"}
    text_content = ""
    try:
        doc_stream = io.BytesIO(file_content)
        doc = docx.Document(doc_stream)
        log_message(f"Processing DOCX '{filename}'.", "info")

        # Extract text from paragraphs
        para_texts = []
        for para in doc.paragraphs:
            para_text = para.text
            para_texts.append(para_text)
            urls.update(find_urls(para_text))
        text_content = '\n'.join(para_texts) # Join with single newline might be better for context
        if text_content:
             parsed_segments.append((text_content, base_meta.copy()))


        # Extract text from tables (basic)
        table_texts = []
        for i, table in enumerate(doc.tables):
             table_str = f"\n[Table {i+1} Content:]\n"
             try:
                  for row in table.rows:
                       row_text = [cell.text for cell in row.cells]
                       table_str += "| " + " | ".join(row_text) + " |\n"
                  table_texts.append(table_str)
             except Exception as table_err:
                  log_message(f"Error reading table {i+1} in {filename}: {table_err}", "warning")
        if table_texts:
            table_meta = base_meta.copy()
            table_meta["content_type"] = "table_text"
            parsed_segments.append(("\n".join(table_texts), table_meta))


        # Handle hyperlinks explicitly
        for rel in doc.part.rels.values():
            if rel.reltype == docx.opc.constants.RELATIONSHIP_TYPE.HYPERLINK and rel.is_external:
                urls.add(rel.target_ref)

        # Image extraction (simplified)
        image_parts = {}
        for rel_id, rel in doc.part.rels.items():
             if "image" in rel.target_ref:
                 try: image_parts[rel_id] = rel.target_part.blob
                 except: pass # Ignore if image part inaccessible

        img_count = 0
        processed_rel_ids = set()
        for shape in doc.inline_shapes:
             try:
                 # Check if it's a picture using XML structure inspection
                 pic = shape._inline.graphic.graphicData.find(f".//{{{docx.oxml.ns.nsmap['pic']}}}pic")
                 if pic is not None:
                     blip = pic.find(f".//{{{docx.oxml.ns.nsmap['a']}}}blip")
                     if blip is not None:
                         embed_id = blip.get(f"{{{docx.oxml.ns.nsmap['r']}}}embed")
                         if embed_id and embed_id in image_parts and embed_id not in processed_rel_ids:
                             image_bytes = image_parts[embed_id]
                             img_meta_filename = f"{filename}_img{img_count}"
                             img_desc_ocr = get_gemini_vision_description_ocr(image_bytes, img_meta_filename)
                             if img_desc_ocr and "No description or text extracted" not in img_desc_ocr:
                                  current_img_meta = img_meta.copy()
                                  current_img_meta["image_ref"] = f"inline_{img_count}"
                                  parsed_segments.append((f"[Analysis of Embedded Image {img_count}]\n{img_desc_ocr}", current_img_meta))
                                  img_count += 1
                                  processed_rel_ids.add(embed_id) # Mark as processed
             except Exception as img_ex:
                 log_message(f"Minor error processing inline shape in {filename}: {img_ex}", "debug")


        log_message(f"Finished processing DOCX '{filename}'. Found URLs: {len(urls)}, Images processed: {img_count}")
    except Exception as e:
        log_message(f"Critical error processing DOCX '{filename}': {e}", "error")
        return [], []
    return parsed_segments, list(urls)

# --- PPTX Parser ---
# Returns list [(slide_text, metadata), (image_analysis_1, metadata), ...]
def parse_pptx(file_content, filename):
    parsed_segments = []
    urls = set()
    try:
        ppt_stream = io.BytesIO(file_content)
        prs = Presentation(ppt_stream)
        log_message(f"Processing PPTX '{filename}': {len(prs.slides)} slides.", "info")

        for i, slide in enumerate(prs.slides):
            slide_num = i + 1
            slide_text_content = ""
            slide_meta = {"source": filename, "slide_number": slide_num, "content_type": "slide_text"}
            slide_urls = set()

            try:
                notes_slide = slide.notes_slide
                if notes_slide and notes_slide.notes_text_frame:
                     notes_text = notes_slide.notes_text_frame.text
                     if notes_text:
                          slide_text_content += f"\n--- Speaker Notes (Slide {slide_num}) ---\n{notes_text}\n"
                          slide_urls.update(find_urls(notes_text))
            except Exception as notes_err:
                 log_message(f"Could not extract speaker notes for slide {slide_num}: {notes_err}", "warning")


            for shape in slide.shapes:
                shape_text = ""
                # Extract text from shapes
                if shape.has_text_frame:
                    try: shape_text = shape.text_frame.text
                    except: pass # Ignore errors getting text
                    if shape_text:
                         slide_text_content += shape_text + "\n"
                         slide_urls.update(find_urls(shape_text))

                # Extract text from tables
                if shape.has_table:
                    table_text = f"\n[Table on Slide {slide_num}:]\n"
                    try:
                        for row in shape.table.rows:
                            row_data = [cell.text_frame.text for cell in row.cells]
                            table_text += "| " + " | ".join(row_data) + " |\n"
                        slide_text_content += table_text
                    except Exception as table_err:
                         log_message(f"Error reading table on slide {slide_num}: {table_err}", "warning")

                # Handle hyperlinks
                try:
                    if shape.click_action.hyperlink and shape.click_action.hyperlink.address:
                        slide_urls.add(shape.click_action.hyperlink.address)
                    if hasattr(shape, 'text_frame'): # Check text runs for hyperlinks
                        for para in shape.text_frame.paragraphs:
                            for run in para.runs:
                                if run.hyperlink and run.hyperlink.address:
                                    slide_urls.add(run.hyperlink.address)
                except: pass # Ignore hyperlink errors

                # Handle images
                if hasattr(shape, 'image'): # shape.shape_type == 13 (Picture) or has image fill
                    try:
                        image = shape.image
                        image_bytes = image.blob
                        img_meta_filename = f"{filename}_slide{slide_num}_img{shape.shape_id}"
                        img_desc_ocr = get_gemini_vision_description_ocr(image_bytes, img_meta_filename)
                        if img_desc_ocr and "No description or text extracted" not in img_desc_ocr:
                             img_meta = {"source": filename, "slide_number": slide_num, "content_type": "embedded_image", "image_ref": f"shape_{shape.shape_id}"}
                             parsed_segments.append((f"[Analysis of Image on Slide {slide_num}]\n{img_desc_ocr}", img_meta))
                    except Exception as img_err:
                         log_message(f"Could not process shape as image on slide {slide_num} (ID {shape.shape_id}): {img_err}", "warning")

            # Add combined text for the slide as one segment
            if slide_text_content.strip():
                 parsed_segments.append((slide_text_content, slide_meta))
            urls.update(slide_urls)

        log_message(f"Finished processing PPTX '{filename}'. Found URLs: {len(urls)}")
    except Exception as e:
        log_message(f"Critical error processing PPTX '{filename}': {e}", "error")
        return [], []
    return parsed_segments, list(urls)

# --- Data File Parsers (CSV/XLSX/JSON) ---
# Returns list [(summary_text, metadata)], data_object
def parse_data_file(file_content, filename):
    """Parses CSV, XLSX, or JSON. Returns list with one summary segment and the data object."""
    text_summary = ""
    data_object = None
    parsed_segments = []
    base_meta = {"source": filename, "content_type": "data_summary"}
    file_type = filename.split('.')[-1].lower()

    try:
        log_message(f"Processing Data file '{filename}'.", "info")
        if file_type == "csv":
            try: df = pd.read_csv(io.BytesIO(file_content), encoding='utf-8')
            except UnicodeDecodeError: df = pd.read_csv(io.BytesIO(file_content), encoding='latin1')
            data_object = df
        elif file_type == "xlsx":
            df = pd.read_excel(io.BytesIO(file_content), engine='openpyxl')
            data_object = df
        elif file_type == "json":
            try: data_object = json.loads(file_content.decode('utf-8'))
            except UnicodeDecodeError: data_object = json.loads(file_content.decode('latin1'))

        if data_object is not None:
            # Store the actual data object in session state
            if "data_frames" not in st.session_state: st.session_state.data_frames = {}
            st.session_state.data_frames[filename] = data_object # Store for later querying

            # Create textual summary for indexing
            text_summary += f"--- Data Summary for {filename} ---\n"
            if isinstance(data_object, pd.DataFrame):
                text_summary += f"Type: Tabular Data (CSV/XLSX)\nShape: {data_object.shape}\n"
                buffer = io.StringIO()
                try: data_object.info(buf=buffer, verbose=True, show_counts=True)
                except Exception: buffer.write("Could not get detailed info.")
                text_summary += "Columns and Data Types:\n" + buffer.getvalue() + "\n"
                text_summary += "First 5 Rows:\n" + data_object.head().to_string() + "\n"
                try:
                    desc = data_object.describe(include='all').to_string()
                    text_summary += f"Basic Statistics:\n{desc}\n"
                except Exception: text_summary += "Basic Statistics: (Could not be generated)\n"
            elif isinstance(data_object, (dict, list)):
                 text_summary += f"Type: JSON Data ({type(data_object)})\n"
                 preview = json.dumps(data_object, indent=2)[:1500] # Limit preview
                 text_summary += f"Data Preview:\n{preview}"
                 if len(preview) >= 1500: text_summary += "...\n" # Indicate truncation
            text_summary += f"--- End Data Summary for {filename} ---\n"

            if text_summary:
                 parsed_segments.append((text_summary, base_meta))
            log_message(f"Successfully parsed and summarized data from '{filename}'.")
        else:
            log_message(f"Could not parse data file '{filename}'.", "warning")

    except Exception as e:
        st.exception(e) # Log full traceback for debugging data parsing errors
        log_message(f"Error processing data file '{filename}': {e}", "error")
        return [], [] # Return empty list on failure

    # Return the list containing the summary segment (if any) and the list of URLs (always empty here)
    return parsed_segments, []

# --- TXT Parser ---
# Returns list [(full_text, metadata)]
def parse_txt(file_content, filename):
    text = ""
    urls = set()
    base_meta = {"source": filename, "content_type": "text"}
    parsed_segments = []
    try:
        log_message(f"Processing TXT file '{filename}'.", "info")
        try: text = file_content.decode("utf-8")
        except UnicodeDecodeError: text = file_content.decode("latin1", errors='ignore')

        if text:
             urls.update(find_urls(text))
             parsed_segments.append((text, base_meta))
             log_message(f"Finished processing TXT '{filename}'. Found URLs: {len(urls)}")
        else:
             log_message(f"TXT file '{filename}' is empty.", "warning")

    except Exception as e:
        log_message(f"Error processing TXT file '{filename}': {e}", "error")
        return [], []
    return parsed_segments, list(urls)

# --- Image File Parser ---
# Returns list [(analysis_text, metadata)]
def parse_image(file_content, filename):
    text = ""
    base_meta = {"source": filename, "content_type": "image_analysis"}
    parsed_segments = []
    try:
        log_message(f"Processing Image file '{filename}'.", "info")
        img_desc_ocr = get_gemini_vision_description_ocr(file_content, filename)
        if img_desc_ocr and "No description or text extracted" not in img_desc_ocr:
             text = f"[Analysis of Image File: {filename}]\n{img_desc_ocr}\n"
             parsed_segments.append((text, base_meta))
             log_message(f"Finished processing image file: {filename}")
        else:
             log_message(f"Image analysis did not return significant content for {filename}.", "warning")

    except Exception as e:
        log_message(f"Error processing image file '{filename}': {e}", "error")
        return [], []
    return parsed_segments, [] # No URLs expected
